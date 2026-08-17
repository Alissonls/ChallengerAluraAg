import os
import json
import csv

SAMPLE_DIR = os.path.join(os.path.dirname(__file__), "data", "sample_docs")
os.makedirs(SAMPLE_DIR, exist_ok=True)

def create_sample_files():
    """Cria documentos de exemplo cobrindo todos os 8 formatos de arquivo e 10 setores corporativos."""
    
    # 1. Markdown - Estratégico (Estrategico_OKRs_Roadmap_2026.md)
    md_content = """# OKRs e Planejamento Estratégico 2026 - Nexus Enterprise

## Objetivos Estratégicos (Q1 - Q4 2026)

### Objetivo 1: Liderança no Mercado de Soluções de Inteligência Artificial Corporativa
- **KR1**: Atingir R$ 50 Milhões em Receita Recorrente Anual (ARR).
- **KR2**: Conquistar 150 novos clientes enterprise na América Latina.
- **KR3**: Manter a taxa de retenção líquida (NDR) acima de 125%.

### Objetivo 2: Excelência Operacional e Infraestrutura Cloud OCI
- **KR1**: Migrar 100% das cargas de trabalho de IA para a Oracle Cloud Infrastructure (OCI).
- **KR2**: Reduzir a latência média de resposta dos agentes RAG para menos de 450ms.
- **KR3**: Atingir índice de disponibilidade (SLA) de 99.99%.

### Diretrizes de Expansão
Expandir as equipes de P&D e Engenharia de IA no primeiro semestre, focando em parcerias estratégicas com ecossistemas de tecnologia.
"""
    with open(os.path.join(SAMPLE_DIR, "Estrategico_OKRs_Roadmap_2026.md"), "w", encoding="utf-8") as f:
        f.write(md_content)

    # 2. JSON - Legal e Compliance (Juridico_LGPD_Politica_Privacidade.json)
    json_legal = {
        "documento": "Política de Privacidade e Conformidade LGPD",
        "versao": "2.4",
        "categoria": "Legal e Compliance",
        "data_vigencia": "2026-01-15",
        "encarregado_dpo": {
            "nome": "Dra. Renata Vasconcelos",
            "email": "dpo@nexusenterprise.com.br",
            "telefone": "+55 11 3000-8800"
        },
        "principios_fundamentais": [
            "Finalidade legítima e específica para cada tratamento de dados pessoais.",
            "Limitação da coleta ao mínimo necessário para execução dos contratos.",
            "Garantia de livre acesso e transparência aos titulares dos dados.",
            "Criptografia ponta a ponta (AES-256) em armazenamento OCI Object Storage."
        ],
        "politica_retencao": {
            "dados_cadastrais": "5 anos após o encerramento do contrato",
            "logs_acesso_agente_ia": "180 dias conforme Marco Civil da Internet",
            "curriculos_rh": "12 meses com consentimento expresso"
        }
    }
    with open(os.path.join(SAMPLE_DIR, "Juridico_LGPD_Politica_Privacidade.json"), "w", encoding="utf-8") as f:
        json.dump(json_legal, f, indent=2, ensure_ascii=False)

    # 3. CSV - Marketing e Comercial (Marketing_Tabela_Precos_Planos.csv)
    csv_rows = [
        ["Plano", "Categoria", "Usuarios_Inclusos", "Preco_Mensal_BRL", "Armazenamento_OCI_GB", "Suporte_SLA"],
        ["Starter AI", "Marketing e Comercial", "25", "1490.00", "50", "24h"],
        ["Business Pro", "Marketing e Comercial", "100", "4990.00", "250", "4h"],
        ["Enterprise Premium", "Marketing e Comercial", "Ilimitado", "12900.00", "1000", "1h Dedicado"],
        ["Custom OCI Cloud", "Marketing e Comercial", "Customizado", "Sob Consulta", "Ilimitado", "Imediato"]
    ]
    with open(os.path.join(SAMPLE_DIR, "Marketing_Tabela_Precos_Planos.csv"), "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerows(csv_rows)

    # 4. HTML - Qualidade (Qualidade_Manual_Auditoria_ISO9001.html)
    html_content = """<!DOCTYPE html>
<html lang="pt-BR">
<head>
    <meta charset="UTF-8">
    <title>Manual de Gestão da Qualidade ISO 9001:2015</title>
</head>
<body>
    <h1>Manual de Gestão da Qualidade (ISO 9001:2015)</h1>
    <p><strong>Setor:</strong> Qualidade e Melhores Práticas</p>
    <p><strong>Código do Documento:</strong> DOC-QUAL-9001-2026</p>

    <h2>1. Objetivo do Sistema de Gestão</h2>
    <p>Garantir o alinhamento contínuo dos processos corporativos, monitorando auditorias internas, não-conformidades e planos de ação corretiva para os produtos corporativos de IA.</p>

    <h2>2. Frequência das Auditorias</h2>
    <ul>
        <li>Auditoria Interna de TI e Infraestrutura Cloud: Trimestral.</li>
        <li>Auditoria de Segurança da Informação e LGPD: Semestral.</li>
        <li>Auditoria Externa de Certificação ISO: Anual.</li>
    </ul>

    <h2>3. Procedimento de Ação Corretiva</h2>
    <p>Sempre que uma falha for identificada em produção, o comitê de qualidade possui até 48 horas para publicar o Relatório de Análise de Causa Raiz (RCA).</p>
</body>
</html>"""
    with open(os.path.join(SAMPLE_DIR, "Qualidade_Manual_Auditoria_ISO9001.html"), "w", encoding="utf-8") as f:
        f.write(html_content)

    # 5. JSON - Dados e Sistemas (Dados_Catalogo_APIs_Integracao.json)
    json_data = {
        "catalogo": "Catálogo de APIs e Microsserviços Corporativos",
        "categoria": "Dados e Sistemas",
        "ambiente_nuvem": "Oracle Cloud Infrastructure (OCI)",
        "endpoints": [
            {
                "nome": "API de Ingestão de Documentos",
                "path": "/api/v1/documents/upload",
                "metodo": "POST",
                "formatos_suportados": ["PDF", "Word", "Excel", "PowerPoint", "Markdown", "CSV", "JSON", "HTML"],
                "autenticacao": "Bearer OAuth2 / OCI IAM Token"
            },
            {
                "nome": "API de RAG & Consulta de Agente",
                "path": "/api/v1/chat/query",
                "metodo": "POST",
                "payload_exemplo": {"query": "Qual a política de reembolso de viagens?", "domain": "Financeiro e Contábil"},
                "tempo_resposta_esperado": "< 500ms"
            }
        ]
    }
    with open(os.path.join(SAMPLE_DIR, "Dados_Catalogo_APIs_Integracao.json"), "w", encoding="utf-8") as f:
        json.dump(json_data, f, indent=2, ensure_ascii=False)

    # 6. Word / DOCX (RH_Politica_Beneficios_HomeOffice.docx)
    try:
        import docx
        doc = docx.Document()
        doc.add_heading('Política de Recursos Humanos: Benefícios e Home Office', 0)
        doc.add_paragraph('Categoria: Recursos Humanos')
        doc.add_heading('1. Regras do Trabalho Híbrido e Home Office', level=1)
        doc.add_paragraph('Todos os colaboradores possuem direito a 3 dias de trabalho remoto por semana e auxílio-home office mensal no valor de R$ 350,00 para despesas de internet e energia.')
        doc.add_heading('2. Benefícios Corporativos', level=1)
        doc.add_paragraph('O plano de saúde é Bradesco Top Nacional sem coparticipação para o titular. O vale refeição/alimentação é disponibilizado via cartão Flash no valor de R$ 1.200,00 mensais.')
        doc.add_heading('3. Solicitação de Férias', level=1)
        doc.add_paragraph('As férias devem ser solicitadas no portal de RH com antecedência mínima de 30 dias após completar o período aquisitivo de 12 meses.')
        doc.save(os.path.join(SAMPLE_DIR, "RH_Politica_Beneficios_HomeOffice.docx"))
    except Exception:
        with open(os.path.join(SAMPLE_DIR, "RH_Politica_Beneficios_HomeOffice.docx"), "w", encoding="utf-8") as f:
            f.write("Política de Recursos Humanos: Benefícios e Home Office\nCategoria: Recursos Humanos\nAuxílio Home Office: R$ 350,00 mensais.\nVale Refeição: R$ 1.200,00 via cartão Flash.\nPlano de Saúde Bradesco Top sem coparticipação.")

    # 7. Excel / XLSX (Financeiro_DRE_Balanço_2025.xlsx)
    try:
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "DRE 2025"
        ws.append(["DRE Sintética", "Categoria", "Q1 2025 (BRL)", "Q2 2025 (BRL)", "Q3 2025 (BRL)", "Q4 2025 (BRL)"])
        ws.append(["Receita Bruta", "Financeiro e Contábil", "12500000", "14200000", "15800000", "18100000"])
        ws.append(["Custos Operacionais OCI Cloud", "Financeiro e Contábil", "-1200000", "-1350000", "-1400000", "-1500000"])
        ws.append(["Despesas com Pessoal (RH)", "Financeiro e Contábil", "-4500000", "-4700000", "-4900000", "-5100000"])
        ws.append(["Lucro Líquido do Exercício", "Financeiro e Contábil", "4800000", "5650000", "6700000", "8200000"])
        wb.save(os.path.join(SAMPLE_DIR, "Financeiro_DRE_Balanco_2025.xlsx"))
    except Exception:
        with open(os.path.join(SAMPLE_DIR, "Financeiro_DRE_Balanco_2025.xlsx"), "w", encoding="utf-8") as f:
            f.write("DRE 2025 | Financeiro e Contábil\nReceita Bruta: R$ 60.600.000,00\nLucro Líquido: R$ 25.350.000,00\nCustos OCI Cloud: R$ 5.450.000,00")

    # 8. PowerPoint / PPTX (Comunicacao_Apresentacao_Onboarding.pptx)
    try:
        import pptx
        prs = pptx.Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        txBox = slide.shapes.add_textbox(0, 0, 100, 100)
        tf = txBox.text_frame
        tf.text = "Guia de Onboarding e Comunicação Interna\nCategoria: Comunicação Interna\nBem-vindo à equipe! Nosso canal oficial de comunicação é o Slack e as reuniões gerais acontecem semanalmente às terças-feiras às 09:00."
        prs.save(os.path.join(SAMPLE_DIR, "Comunicacao_Apresentacao_Onboarding.pptx"))
    except Exception:
        with open(os.path.join(SAMPLE_DIR, "Comunicacao_Apresentacao_Onboarding.pptx"), "w", encoding="utf-8") as f:
            f.write("Guia de Onboarding e Comunicação Interna\nCategoria: Comunicação Interna\nReunião Semanal: Terças às 09h.\nCanal de Comunicação Oficial: Slack.")

    # 9. PDF (Operacional_Manual_Procedimentos_TI.pdf)
    try:
        import pypdf
        writer = pypdf.PdfWriter()
        page = writer.add_blank_page(width=612, height=792)
        with open(os.path.join(SAMPLE_DIR, "Operacional_Manual_Procedimentos_TI.pdf"), "w", encoding="utf-8") as f:
            f.write("Manual de Procedimentos Operacionais de TI\nCategoria: Operacional\nCódigo: MAN-TI-001\nProcedimento em caso de falha no servidor: Notificar equipe DevOps em #ti-incidentes no Slack. O tempo máximo de resposta para severidade 1 é 15 minutos.")
    except Exception:
        with open(os.path.join(SAMPLE_DIR, "Operacional_Manual_Procedimentos_TI.pdf"), "w", encoding="utf-8") as f:
            f.write("Manual Operacional TI | Categoria: Operacional | Tempo de resposta para incidentes críticos: 15 minutos.")

    # 10. PDF (Pesquisa_Business_Case_AI.pdf)
    with open(os.path.join(SAMPLE_DIR, "Pesquisa_Business_Case_AI.pdf"), "w", encoding="utf-8") as f:
        f.write("Business Case: Agente Corporativo de IA em Oracle Cloud (OCI)\nCategoria: Pesquisa e Desenvolvimento\nResumo: Estudo de viabilidade de agentes RAG corporativos demonstrando ganho de 35% de produtividade na busca por informações internas.")

if __name__ == "__main__":
    create_sample_files()
    print("Documentos corporativos de exemplo gerados com sucesso!")
