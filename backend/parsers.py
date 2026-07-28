import os
import json
import csv
import io
import re
from typing import Dict, Any, List

def parse_pdf(file_bytes: bytes, filename: str) -> str:
    """Extrai texto de um arquivo PDF usando a biblioteca pypdf."""
    text_chunks = []
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        for idx, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text_chunks.append(f"[Página {idx + 1}]\n{page_text.strip()}")
    except Exception as e:
        # Fallback para extração direta de caracteres se a biblioteca falhar ou o fluxo for parcial
        printable = re.sub(r'[^\x20-\x7E\n\r\t]', ' ', file_bytes.decode('utf-8', errors='ignore'))
        clean_lines = [line.strip() for line in printable.split('\n') if len(line.strip()) > 3]
        text_chunks.append("\n".join(clean_lines[:100]))
    
    return "\n\n".join(text_chunks) if text_chunks else "Conteúdo PDF sem texto extraível."

def parse_docx(file_bytes: bytes, filename: str) -> str:
    """Extrai texto de um arquivo Word .docx usando a biblioteca python-docx."""
    try:
        import docx
        doc = docx.Document(io.BytesIO(file_bytes))
        full_text = []
        for p in doc.paragraphs:
            if p.text.strip():
                full_text.append(p.text.strip())
        for table in doc.tables:
            for row in table.rows:
                row_str = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                if row_str:
                    full_text.append(f"[Tabela] {row_str}")
        return "\n".join(full_text)
    except Exception as e:
        # Fallback para leitura direta em UTF-8
        return file_bytes.decode('utf-8', errors='ignore')

def parse_excel(file_bytes: bytes, filename: str) -> str:
    """Extrai texto de um arquivo Excel .xlsx/.xls usando openpyxl ou pandas."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        sheets_text = []
        for sheetname in wb.sheetnames:
            ws = wb[sheetname]
            sheet_rows = []
            for row in ws.iter_rows(values_only=True):
                row_vals = [str(val) for val in row if val is not None]
                if row_vals:
                    sheet_rows.append(" | ".join(row_vals))
            if sheet_rows:
                sheets_text.append(f"=== Planilha: {sheetname} ===\n" + "\n".join(sheet_rows))
        return "\n\n".join(sheets_text)
    except Exception as e:
        try:
            import pandas as pd
            df_dict = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
            res = []
            for name, df in df_dict.items():
                res.append(f"=== Planilha: {name} ===\n" + df.to_string())
            return "\n\n".join(res)
        except Exception:
            return file_bytes.decode('utf-8', errors='ignore')

def parse_pptx(file_bytes: bytes, filename: str) -> str:
    """Extrai texto de uma apresentação PowerPoint .pptx usando python-pptx."""
    try:
        import pptx
        prs = pptx.Presentation(io.BytesIO(file_bytes))
        slides_text = []
        for idx, slide in enumerate(prs.slides):
            slide_lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_lines.append(shape.text.strip())
            if slide_lines:
                slides_text.append(f"[Slide {idx + 1}]\n" + "\n".join(slide_lines))
        return "\n\n".join(slides_text)
    except Exception as e:
        return file_bytes.decode('utf-8', errors='ignore')

def parse_markdown(file_bytes: bytes, filename: str) -> str:
    """Extrai texto de um arquivo Markdown."""
    return file_bytes.decode('utf-8', errors='ignore')

def parse_csv(file_bytes: bytes, filename: str) -> str:
    """Extrai texto de um arquivo CSV tabulado."""
    decoded = file_bytes.decode('utf-8', errors='ignore')
    reader = csv.reader(io.StringIO(decoded))
    lines = []
    for row in reader:
        if row:
            lines.append(" | ".join(row))
    return "\n".join(lines)

def parse_json(file_bytes: bytes, filename: str) -> str:
    """Extrai texto de um arquivo JSON e formata a estrutura."""
    try:
        decoded = file_bytes.decode('utf-8', errors='ignore')
        data = json.loads(decoded)
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception:
        return file_bytes.decode('utf-8', errors='ignore')

def parse_html(file_bytes: bytes, filename: str) -> str:
    """Extrai texto limpo de uma página HTML usando BeautifulSoup."""
    decoded = file_bytes.decode('utf-8', errors='ignore')
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(decoded, 'html.parser')
        # Remove elementos de script e estilo
        for script in soup(["script", "style", "head"]):
            script.extract()
        text = soup.get_text(separator='\n')
        # Divide em linhas e remove espaços extras
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)
        return text
    except Exception:
        return re.sub('<[^<]+?>', '', decoded)

def parse_document(file_bytes: bytes, filename: str) -> Dict[str, Any]:
    """
    Roteador principal de extração de documentos com base na extensão do arquivo.
    Suporta PDF, Word, Excel, PowerPoint, Markdown, CSV, JSON e HTML.
    """
    ext = os.path.splitext(filename)[1].lower()
    
    if ext == '.pdf':
        content = parse_pdf(file_bytes, filename)
        doc_type = 'PDF'
    elif ext in ['.docx', '.doc']:
        content = parse_docx(file_bytes, filename)
        doc_type = 'Word'
    elif ext in ['.xlsx', '.xls']:
        content = parse_excel(file_bytes, filename)
        doc_type = 'Excel'
    elif ext in ['.pptx', '.ppt']:
        content = parse_pptx(file_bytes, filename)
        doc_type = 'PowerPoint'
    elif ext in ['.md', '.markdown']:
        content = parse_markdown(file_bytes, filename)
        doc_type = 'Markdown'
    elif ext == '.csv':
        content = parse_csv(file_bytes, filename)
        doc_type = 'CSV'
    elif ext == '.json':
        content = parse_json(file_bytes, filename)
        doc_type = 'JSON'
    elif ext in ['.html', '.htm']:
        content = parse_html(file_bytes, filename)
        doc_type = 'HTML'
    else:
        content = file_bytes.decode('utf-8', errors='ignore')
        doc_type = 'Text'
        
    return {
        "filename": filename,
        "extension": ext,
        "doc_type": doc_type,
        "content": content,
        "char_count": len(content),
        "word_count": len(content.split())
    }

